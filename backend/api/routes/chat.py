"""
Chat API Router
- Orchestrates multi-provider routing and RAG pipeline
- Integrates service detection for improved prompt engineering
- Uses loguru for logging
"""
from fastapi import APIRouter, Depends
from sqlalchemy.ext.asyncio import AsyncSession
from backend.api.schemas import ChatRequest, ChatResponse
from backend.services.router import QueryRouter
from backend.models.database import get_db
from backend.utils.service_detection import detect_service
from backend.utils.source_validator import (
    validate_source, 
    get_valid_sources,
    get_correct_source_for_service
)
from loguru import logger
import uuid
import base64
import os
import io
from fastapi import UploadFile, File, Form
from typing import Optional
from backend.services.llm_service import llm_service

router = APIRouter()

@router.post("/", response_model=ChatResponse)
async def chat_query(request: ChatRequest, db: AsyncSession = Depends(get_db)):
    """Handles user queries by routing between documents (RAG) and Live Cloud APIs."""
    try:
        query_router = QueryRouter(db)
        result = await query_router.route(request.query, request.selected_source, request.selected_db)
        
        from backend.services.llm_service import LLMService
        llm_service = LLMService()
        
        context = ""
        sources = []
        
        # Detect target service for prompt specialization
        primary_service = detect_service(request.query)

        if result["source_type"] == "docs":
            # Extract clean sources and format context
            context_parts = []
            chunks = result.get("data") or result.get("chunks") or result.get("results") or []
            
            for chunk in chunks:
                if isinstance(chunk, dict):
                    src = chunk.get("source") or chunk.get("metadata", {}).get("source", "")
                    content = chunk.get("content", "")
                    chunk_idx = chunk.get("metadata", {}).get("chunk_index", "?")
                elif hasattr(chunk, "metadata"):
                    src = chunk.metadata.get("source", "")
                    content = chunk.page_content if hasattr(chunk, "page_content") else str(chunk)
                    chunk_idx = chunk.metadata.get("chunk_index", "?")
                else:
                    src = ""
                    content = str(chunk)
                    chunk_idx = "?"
                
                if src and src not in sources:
                    sources.append(src)
                
                if content:
                    context_parts.append(f"SOURCE: {src} (Chunk {chunk_idx})\nCONTENT: {content}")
            
            # Validate all sources before passing to LLM
            valid = get_valid_sources()
            validated_sources = []
            for s in sources:
                if any(v.lower() in s.lower() for v in valid):
                    validated_sources.append(s)
            
            if not validated_sources and primary_service:
                validated_sources = [get_correct_source_for_service(primary_service)]
            
            sources = validated_sources
            context = "\n\n".join(context_parts)
            
        elif result["source_type"] == "api":
            context = str(result["data"])
            sources = [f"{result.get('provider', 'AWS').upper()} API: {result.get('service', '').upper()}"]
            
        from backend.services.learning_service import LearningService
        learning_service = LearningService(db)
        user_id = "default_user" 
        insights = await learning_service.get_user_insights(user_id)

        # Generate response using the improved RAG prompt
        answer = await llm_service.generate_response(
            request.query, 
            context, 
            sources,
            service=primary_service,
            learned_insights=insights
        )
        
        logger.info(f"Query processed successfully. Source: {result['source_type']}")
        
        return ChatResponse(
            answer=answer,
            source_type=result["source_type"],
            source_details=result.get("data", []) or result.get("chunks", []) or result.get("results", []),
            conversation_id=request.conversation_id or str(uuid.uuid4().hex)
        )
    except Exception as e:
        logger.error(f"Chat execution failed: {e}")
        return ChatResponse(
            answer=f"I encountered an error processing your request: {str(e)}",
            source_type="error",
            source_details=[],
            conversation_id=request.conversation_id or str(uuid.uuid4().hex)
        )

@router.post("/with-file")
async def chat_with_file(
    message: str = Form(...),
    file: Optional[UploadFile] = File(None)
):
    """
    Chat endpoint that accepts an optional 
    file (image, PDF, txt, docx, pptx etc)
    and answers questions about it.
    """
    file_context = ""
    is_image = False
    image_base64 = None

    if file:
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()
        content = await file.read()

        # ── IMAGES ──────────────────────────────
        if ext in (".png", ".jpg", ".jpeg",
                   ".gif", ".bmp", ".webp", ".tiff"):
            is_image = True
            
            # Convert to base64 for vision model
            image_base64 = base64.b64encode(
                content
            ).decode("utf-8")
            
            # Also run OCR to extract any text
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(io.BytesIO(content))
                ocr_text = pytesseract.image_to_string(img)
                if ocr_text.strip():
                    file_context = (
                        f"Text extracted from image "
                        f"via OCR:\n{ocr_text}"
                    )
            except Exception as e:
                logger.warning(f"OCR failed: {e}")

        # ── PDFs ────────────────────────────────
        elif ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(
                    io.BytesIO(content)
                ) as pdf:
                    pages = []
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            pages.append(
                                f"Page {i+1}:\n{text}"
                            )
                file_context = "\n\n".join(pages)
            except Exception as e:
                logger.error(f"PDF read failed: {e}")

        # ── WORD DOCS ───────────────────────────
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(io.BytesIO(content))
                file_context = "\n".join(
                    p.text for p in doc.paragraphs 
                    if p.text.strip()
                )
            except Exception as e:
                logger.error(f"DOCX read failed: {e}")

        # ── PPTX ────────────────────────────────
        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in (
                                shape.text_frame.paragraphs
                            ):
                                if para.text.strip():
                                    texts.append(para.text)
                    if texts:
                        slides.append(
                            f"Slide {i+1}:\n"
                            + "\n".join(texts)
                        )
                file_context = "\n\n".join(slides)
            except Exception as e:
                logger.error(f"PPTX read failed: {e}")

        # ── TEXT / JSON / MD / CSV ───────────────
        elif ext in (".txt", ".md", ".json",
                     ".csv", ".yaml", ".xml",
                     ".py", ".js", ".sql"):
            try:
                file_context = content.decode(
                    "utf-8", errors="ignore"
                )
            except Exception as e:
                logger.error(f"Text read failed: {e}")

        else:
            file_context = (
                f"Unsupported file type: {ext}. "
                f"Supported: images, PDF, DOCX, "
                f"PPTX, TXT, JSON, CSV, MD"
            )

    # ── BUILD PROMPT WITH FILE CONTEXT ──────────
    if is_image and image_base64:
        # Use vision model (llava) for images
        response = await llm_service.analyze_image(
            message=message,
            image_base64=image_base64,
            ocr_context=file_context
        )
    elif file_context:
        # Use text model with file content as context
        prompt = f"""The user uploaded a file and 
is asking a question about it.

FILE CONTENT:
{file_context[:8000]}

USER QUESTION:
{message}

Answer the question based on the file content above.
If it is an AWS architecture diagram or document,
explain it clearly for someone learning AWS.
Be specific and reference actual content from the file."""

        response = await llm_service.generate_response(
            query=prompt,
            context="",
            sources=[file.filename if file else "Uploaded File"],
            service="general"
        )
    else:
        # No file — normal RAG chat
        from backend.api.schemas import ChatRequest
        from backend.models.database import get_db
        # This is a bit tricky because we'd need to mock DB access or refactor
        # For now, let's keep it simple as requested
        response = "Please provide a file or ask a question about your knowledge base."

    return {
        "answer": response,
        "has_file": file is not None,
        "file_type": (
            os.path.splitext(file.filename)[1] 
            if file else None
        ),
        "file_name": file.filename if file else None
    }
