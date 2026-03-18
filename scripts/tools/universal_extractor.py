import os
import json
import zipfile
import py7zr
import shutil
import time
import boto3
import pandas as pd
import requests as req
from pathlib import Path
from datetime import datetime
from loguru import logger
from concurrent.futures import ThreadPoolExecutor, as_completed

# Image processing
from PIL import Image
import pytesseract

# PDF processing
import pdfplumber
from pdf2image import convert_from_path

# Office documents
import docx
from pptx import Presentation
import openpyxl

# Web content
import html2text
import markdown

S3_BUCKET = "cloud-intelligence-rag-store"
BOT_URL = "http://localhost:8000"
OUTPUT_DIR = "data/extracted"
s3 = boto3.client("s3", region_name="us-east-1")

# Track everything
REPORT = {
    "started": datetime.now().isoformat(),
    "files_found": [],
    "extracted": [],
    "failed": [],
    "s3_uploaded": [],
    "bot_ingested": []
}

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(f"{OUTPUT_DIR}/text", exist_ok=True)
os.makedirs(f"{OUTPUT_DIR}/images", exist_ok=True)
os.makedirs(f"{OUTPUT_DIR}/tables", exist_ok=True)
os.makedirs(f"{OUTPUT_DIR}/raw", exist_ok=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# EXTRACTORS PER FILE TYPE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def extract_pdf(filepath: str) -> dict:
    """Extract text + tables + images from PDF"""
    results = {"text": "", "tables": [], "images": []}
    try:
        # Extract text and tables
        with pdfplumber.open(filepath) as pdf:
            all_text = []
            for i, page in enumerate(pdf.pages):
                # Text
                text = page.extract_text()
                if text:
                    all_text.append(
                        f"--- Page {i+1} ---\n{text}"
                    )
                # Tables
                tables = page.extract_tables()
                for t_idx, table in enumerate(tables):
                    if table:
                        df = pd.DataFrame(table)
                        table_text = df.to_string()
                        results["tables"].append(
                            f"Page {i+1} Table {t_idx+1}:"
                            f"\n{table_text}"
                        )
            results["text"] = "\n\n".join(all_text)

        # Extract images using OCR
        try:
            images = convert_from_path(filepath, dpi=200)
            for i, img in enumerate(images):
                img_path = (
                    f"{OUTPUT_DIR}/images/"
                    f"{Path(filepath).stem}_page{i+1}.png"
                )
                img.save(img_path)
                
                # OCR on the image
                ocr_text = pytesseract.image_to_string(
                    img, config="--psm 6"
                )
                if ocr_text.strip():
                    results["images"].append(
                        f"Page {i+1} OCR:\n{ocr_text}"
                    )
                    
        except Exception as e:
            logger.warning(f"PDF image extract failed: {e}")

        logger.info(
            f"PDF extracted: {Path(filepath).name} "
            f"({len(results['text'])} chars, "
            f"{len(results['tables'])} tables)"
        )
    except Exception as e:
        logger.error(f"PDF failed {filepath}: {e}")
    return results


def extract_image(filepath: str) -> dict:
    """Extract text from image using OCR"""
    results = {"text": "", "description": ""}
    try:
        img = Image.open(filepath)
        
        # Convert to RGB if needed
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        
        # OCR extraction
        ocr_text = pytesseract.image_to_string(
            img, config="--psm 6"
        )
        
        # Get image metadata
        width, height = img.size
        results["description"] = (
            f"Image: {Path(filepath).name} "
            f"({width}x{height} pixels, {img.mode} mode)"
        )
        results["text"] = ocr_text
        
        logger.info(
            f"Image OCR: {Path(filepath).name} "
            f"({len(ocr_text)} chars extracted)"
        )
    except Exception as e:
        logger.error(f"Image failed {filepath}: {e}")
    return results


def extract_docx(filepath: str) -> dict:
    """Extract text + tables from Word document"""
    results = {"text": "", "tables": []}
    try:
        doc = docx.Document(filepath)
        
        # Extract paragraphs
        paragraphs = [
            p.text for p in doc.paragraphs 
            if p.text.strip()
        ]
        results["text"] = "\n".join(paragraphs)
        
        # Extract tables
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(" | ".join(cells))
            results["tables"].append(
                f"Table {i+1}:\n" + "\n".join(rows)
            )
        
        logger.info(
            f"DOCX extracted: {Path(filepath).name}"
        )
    except Exception as e:
        logger.error(f"DOCX failed {filepath}: {e}")
    return results


def extract_pptx(filepath: str) -> dict:
    """Extract text + images from PowerPoint"""
    results = {"text": "", "images": []}
    try:
        prs = Presentation(filepath)
        slides_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = [f"--- Slide {i+1} ---"]
            
            for shape in slide.shapes:
                # Text
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_content.append(para.text)
                
                # Images in slides
                if shape.shape_type == 13:  # Picture
                    try:
                        img_bytes = shape.image.blob
                        img_path = (
                            f"{OUTPUT_DIR}/images/"
                            f"{Path(filepath).stem}"
                            f"_slide{i+1}_{shape.shape_id}.png"
                        )
                        with open(img_path, "wb") as f:
                            f.write(img_bytes)
                        
                        # OCR the extracted image
                        img = Image.open(img_path)
                        ocr = pytesseract.image_to_string(img)
                        if ocr.strip():
                            slide_content.append(
                                f"[Image text: {ocr.strip()}]"
                            )
                    except Exception as e:
                        logger.warning(
                            f"PPTX image failed: {e}"
                        )
            
            slides_text.append("\n".join(slide_content))
        
        results["text"] = "\n\n".join(slides_text)
        logger.info(
            f"PPTX extracted: {Path(filepath).name} "
            f"({len(prs.slides)} slides)"
        )
    except Exception as e:
        logger.error(f"PPTX failed {filepath}: {e}")
    return results


def extract_xlsx(filepath: str) -> dict:
    """Extract all sheets from Excel"""
    results = {"text": ""}
    try:
        xl = pd.ExcelFile(filepath)
        all_sheets = []
        
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(
                filepath, sheet_name=sheet_name
            )
            sheet_text = (
                f"Sheet: {sheet_name}\n"
                f"{df.to_string(index=False)}"
            )
            all_sheets.append(sheet_text)
        
        results["text"] = "\n\n".join(all_sheets)
        logger.info(
            f"XLSX extracted: {Path(filepath).name} "
            f"({len(xl.sheet_names)} sheets)"
        )
    except Exception as e:
        logger.error(f"XLSX failed {filepath}: {e}")
    return results


def extract_csv(filepath: str) -> dict:
    """Extract CSV as text"""
    results = {"text": ""}
    try:
        df = pd.read_csv(filepath)
        results["text"] = df.to_string(index=False)
        logger.info(f"CSV extracted: {Path(filepath).name}")
    except Exception as e:
        logger.error(f"CSV failed {filepath}: {e}")
    return results


def extract_html(filepath: str) -> dict:
    """Extract clean text from HTML"""
    results = {"text": ""}
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            html_content = f.read()
        h = html2text.HTML2Text()
        h.ignore_links = False
        results["text"] = h.handle(html_content)
        logger.info(
            f"HTML extracted: {Path(filepath).name}"
        )
    except Exception as e:
        logger.error(f"HTML failed {filepath}: {e}")
    return results


def extract_text_file(filepath: str) -> dict:
    """Read plain text, markdown, json files"""
    results = {"text": ""}
    try:
        with open(filepath, "r", 
                  encoding="utf-8", errors="ignore") as f:
            results["text"] = f.read()
        logger.info(
            f"Text extracted: {Path(filepath).name}"
        )
    except Exception as e:
        logger.error(f"Text failed {filepath}: {e}")
    return results


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN EXTRACTOR — handles any file type
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def extract_file(filepath: str) -> str:
    """Route file to correct extractor, return text"""
    ext = Path(filepath).suffix.lower()
    name = Path(filepath).name
    
    REPORT["files_found"].append(name)
    
    try:
        if ext == ".pdf":
            data = extract_pdf(filepath)
            parts = [data["text"]]
            parts.extend(data.get("tables", []))
            parts.extend(data.get("images", []))
            text = "\n\n".join(filter(None, parts))

        elif ext in (".png", ".jpg", ".jpeg",
                     ".gif", ".bmp", ".tiff", ".webp"):
            data = extract_image(filepath)
            text = (
                f"{data['description']}\n\n"
                f"Extracted Text:\n{data['text']}"
            )

        elif ext == ".docx":
            data = extract_docx(filepath)
            text = data["text"]
            if data["tables"]:
                text += "\n\n" + "\n\n".join(
                    data["tables"]
                )

        elif ext in (".pptx", ".ppt"):
            data = extract_pptx(filepath)
            text = data["text"]

        elif ext in (".xlsx", ".xls"):
            data = extract_xlsx(filepath)
            text = data["text"]

        elif ext == ".csv":
            data = extract_csv(filepath)
            text = data["text"]

        elif ext in (".html", ".htm"):
            data = extract_html(filepath)
            text = data["text"]

        elif ext in (".txt", ".md", ".json",
                     ".yaml", ".yml", ".xml",
                     ".py", ".js", ".ts",
                     ".sh", ".sql", ".log"):
            data = extract_text_file(filepath)
            text = data["text"]

        else:
            logger.warning(f"Unsupported type: {ext}")
            return ""

        return text.strip()

    except Exception as e:
        logger.error(f"Extract failed {filepath}: {e}")
        REPORT["failed"].append(name)
        return ""


def save_as_txt(text: str, 
                original_path: str, 
                source_label: str) -> str:
    """Save extracted text as .txt file"""
    original_name = Path(original_path).name
    # Replace dots with underscores expect for the last one if we want to keep it?
    # Actually just append .txt to the whole filename
    out_name = f"{original_name}.txt"
    out_path = f"{OUTPUT_DIR}/text/{out_name}"
    
    header = (
        f"SOURCE FILE: {original_name}\n"
        f"SOURCE LABEL: {source_label}\n"
        f"EXTRACTED: {datetime.now().isoformat()}\n"
        f"{'='*60}\n\n"
    )
    
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(header + text)
    
    REPORT["extracted"].append(out_name)
    return out_path


def upload_to_s3(local_path: str, category: str):
    """Upload extracted file to S3"""
    try:
        fname = Path(local_path).name
        key = f"documents/extracted/{category}/{fname}"
        s3.upload_file(local_path, S3_BUCKET, key)
        REPORT["s3_uploaded"].append(fname)
        logger.info(f"☁️  S3: {key}")
    except Exception as e:
        logger.error(f"S3 failed {local_path}: {e}")


def ingest_into_bot(local_path: str):
    """Send extracted text file to RAG bot"""
    try:
        fname = Path(local_path).name
        with open(local_path, "rb") as f:
            r = req.post(
                f"{BOT_URL}/api/v1/documents/upload",
                files={"file": (fname, f, "text/plain")},
                timeout=120
            )
        if r.status_code == 200:
            REPORT["bot_ingested"].append(fname)
            logger.info(f"🤖 BOT: {fname}")
        else:
            logger.error(
                f"BOT failed {fname}: {r.status_code}"
            )
    except Exception as e:
        logger.error(f"BOT error {local_path}: {e}")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PROCESS ZIP OR FOLDER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def process_zip(zip_path: str):
    """Extract ZIP and process all contents"""
    logger.info(f"📦 Opening ZIP: {zip_path}")
    
    extract_to = f"data/extracted/raw/{Path(zip_path).stem}"
    os.makedirs(extract_to, exist_ok=True)
    
    with zipfile.ZipFile(zip_path, "r") as z:
        z.extractall(extract_to)
        logger.info(
            f"ZIP extracted to: {extract_to} "
            f"({len(z.namelist())} files)"
        )
    
    process_folder(extract_to, 
                   source_label=Path(zip_path).name)


def process_single_file(filepath: str, root: str, folder_path: str, source_label: str):
    """Worker function for parallel processing"""
    filename = Path(filepath).name
    rel_path = os.path.relpath(filepath, folder_path)
    
    logger.info(f"Processing: {rel_path}")
    
    # Extract content
    text = extract_file(filepath)
    
    if not text or len(text) < 50:
        logger.warning(f"Skipped (no content): {filename}")
        return
    
    # Determine category from subfolder name
    subfolder = Path(root).name
    category = subfolder if subfolder != Path(folder_path).name else "root"
    
    # Save as .txt
    out_path = save_as_txt(
        text, filepath,
        f"{source_label}/{rel_path}"
    )
    
    # Upload to S3
    upload_to_s3(out_path, category)
    
    # Ingest into RAG bot
    ingest_into_bot(out_path)


def process_folder(folder_path: str, 
                   source_label: str = "folder"):
    """
    Recursively process ALL files in folder
    including all subfolders using parallel workers
    """
    logger.info(f"📂 Processing folder: {folder_path}")
    
    all_tasks = []
    
    # Walk through ALL subfolders recursively to collect files
    for root, dirs, files in os.walk(folder_path):
        # Skip hidden folders
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        
        for filename in files:
            # Skip hidden files and system files
            if filename.startswith(".") or filename in [
                "__MACOSX", "Thumbs.db", "desktop.ini"
            ]:
                continue
            
            filepath = os.path.join(root, filename)
            all_tasks.append((filepath, root))

    # Process in parallel
    if all_tasks:
        max_workers = min(10, len(all_tasks))
        logger.info(f"🚀 Starting parallel extraction with {max_workers} workers")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = [
                executor.submit(process_single_file, fp, r, folder_path, source_label) 
                for fp, r in all_tasks
            ]
            for future in as_completed(futures):
                try:
                    future.result()
                except Exception as e:
                    logger.error(f"Worker thread failed: {e}")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ENTRY POINT — Change path here
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage:")
        print("  For ZIP:    python3 universal_extractor.py myfile.zip")
        print("  For folder: python3 universal_extractor.py /path/to/folder")
        sys.exit(1)
    
    input_path = sys.argv[1]
    
    if not os.path.exists(input_path):
        print(f"❌ Path not found: {input_path}")
        sys.exit(1)
    
    # Route to correct processor
    if input_path.endswith(".zip"):
        process_zip(input_path)
    elif os.path.isdir(input_path):
        process_folder(input_path)
    else:
        # Single file
        text = extract_file(input_path)
        if text:
            out = save_as_txt(
                text, input_path, "single_file"
            )
            upload_to_s3(out, "single")
            ingest_into_bot(out)
    
    # Save final report
    REPORT["finished"] = datetime.now().isoformat()
    REPORT["summary"] = {
        "files_found": len(REPORT["files_found"]),
        "extracted": len(REPORT["extracted"]),
        "failed": len(REPORT["failed"]),
        "s3_uploaded": len(REPORT["s3_uploaded"]),
        "bot_ingested": len(REPORT["bot_ingested"])
    }
    
    report_path = (
        f"{OUTPUT_DIR}/extraction_report_"
        f"{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        f".json"
    )
    with open(report_path, "w") as f:
        json.dump(REPORT, f, indent=2)
    
    # Upload report to S3
    s3.upload_file(
        report_path, S3_BUCKET,
        f"registry/extraction_reports/"
        f"{Path(report_path).name}"
    )
    
    print("\n" + "="*50)
    print("EXTRACTION COMPLETE")
    print("="*50)
    s = REPORT["summary"]
    print(f"Files Found:   {s['files_found']}")
    print(f"Extracted:     {s['extracted']}")
    print(f"Failed:        {s['failed']}")
    print(f"S3 Uploaded:   {s['s3_uploaded']}")
    print(f"Bot Ingested:  {s['bot_ingested']}")
    print(f"Report:        {report_path}")
