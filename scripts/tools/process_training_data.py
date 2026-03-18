import os
import io
import json
import boto3
import base64
import httpx
import asyncio
import numpy as np
import tempfile
from pathlib import Path
from datetime import datetime, timezone
from concurrent.futures import ThreadPoolExecutor, as_completed
from sentence_transformers import SentenceTransformer

# Import existing app services
import sys
sys.path.insert(0, '.')
from backend.core.config import settings

S3_BUCKET = 'cloud-intelligence-rag-store'
EXTRACT_DIR = 'data/extracted/training'
MAX_FILES_TO_PROCESS = 10 # Added safety limit for demo
s3 = boto3.client('s3', region_name='us-east-1')
embedder = SentenceTransformer('all-MiniLM-L6-v2')

REPORT = {
    'started': datetime.now().isoformat(),
    'files_processed': [],
    'files_failed': [],
    's3_uploaded': [],
    'total_chunks': 0
}

# ── TEXT EXTRACTION PER FILE TYPE ──────────────

def extract_text(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    
    try:
        if ext == '.pdf':
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text()
                    if t:
                        pages.append(
                            f'Page {i+1}:\n{t}'
                        )
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            import pandas as pd
                            df = pd.DataFrame(table)
                            pages.append(df.to_string())
            return '\n\n'.join(pages)

        elif ext in ('.png','.jpg','.jpeg',
                     '.gif','.bmp','.webp','.tiff'):
            # OCR for text extraction
            from PIL import Image
            import pytesseract
            img = Image.open(filepath)
            if img.mode not in ('RGB','L'):
                img = img.convert('RGB')
            ocr = pytesseract.image_to_string(
                img, config='--psm 6'
            )
            
            # Also use llava for visual description
            with open(filepath, 'rb') as f:
                img_b64 = base64.b64encode(
                    f.read()
                ).decode('utf-8')
            
            try:
                # Fix: Ollama URL usually needs to be formatted
                # Replacing /v1 if it exists as /api/generate is the raw endpoint
                ollama_url = f"{settings.OLLAMA_BASE_URL.replace('/v1', '')}/api/generate"
                r = httpx.post(
                    ollama_url,
                    json={
                        'model': 'llava',
                        'prompt': (
                            'Describe this image in '
                            'detail. Extract all text, '
                            'labels, arrows, components '
                            'and explain what this '
                            'diagram or image shows. '
                            'Be thorough and specific.'
                        ),
                        'images': [img_b64],
                        'stream': False,
                        'options': {'temperature': 0.1}
                    },
                    timeout=120
                )
                vision_text = r.json().get(
                    'response', ''
                ) if r.status_code == 200 else ''
            except Exception as e:
                print(f"Ollama/Llava failed: {e}")
                vision_text = ''
            
            combined = []
            if ocr.strip():
                combined.append(
                    f'OCR Text:\n{ocr.strip()}'
                )
            if vision_text.strip():
                combined.append(
                    f'Visual Description:\n'
                    f'{vision_text.strip()}'
                )
            return '\n\n'.join(combined)

        elif ext == '.docx':
            import docx
            doc = docx.Document(filepath)
            parts = [
                p.text for p in doc.paragraphs 
                if p.text.strip()
            ]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(
                        ' | '.join(
                            c.text for c in row.cells
                        )
                    )
            return '\n'.join(parts)

        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(filepath)
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in (
                            shape.text_frame.paragraphs
                        ):
                            if p.text.strip():
                                texts.append(p.text)
                if texts:
                    slides.append(
                        f'Slide {i+1}:\n'
                        + '\n'.join(texts)
                    )
            return '\n\n'.join(slides)

        elif ext in ('.xlsx', '.xls'):
            import pandas as pd
            xl = pd.ExcelFile(filepath)
            sheets = []
            for name in xl.sheet_names:
                df = pd.read_excel(
                    filepath, sheet_name=name
                )
                sheets.append(
                    f'Sheet {name}:\n'
                    f'{df.to_string(index=False)}'
                )
            return '\n\n'.join(sheets)

        elif ext == '.csv':
            import pandas as pd
            df = pd.read_csv(filepath)
            return df.to_string(index=False)

        elif ext in ('.html', '.htm'):
            import html2text
            with open(filepath, 'r',
                      errors='ignore') as f:
                return html2text.html2text(f.read())

        elif ext in ('.txt','.md','.json',
                     '.yaml','.yml','.xml',
                     '.py','.js','.ts',
                     '.sh','.sql','.log'):
            with open(filepath, 'r',
                      errors='ignore') as f:
                return f.read()
        
        else:
            # print(f'Skipping unsupported: {ext}')
            return ''
            
    except Exception as e:
        print(f'Extract error {filepath}: {e}')
        return ''


def chunk_text(text: str, 
               filename: str,
               chunk_size: int = 500,
               overlap: int = 50) -> list:
    words = text.split()
    chunks = []
    if not words:
        return []
    i = 0
    chunk_idx = 0
    while i < len(words):
        chunk_words = words[i:i+chunk_size]
        chunk_text = ' '.join(chunk_words)
        chunks.append({
            'text': chunk_text,
            'metadata': {
                'source': filename,
                'chunk_index': chunk_idx,
                'total_chars': len(chunk_text),
                'doc_type': 'training_data'
            }
        })
        i += chunk_size - overlap
        chunk_idx += 1
    return chunks


def vectorize_and_upload(
    filename: str,
    filepath: str,
    service: str,
    chunks: list
):
    if not chunks:
        return 0
    
    # Vectorize all chunks at once (fast)
    texts = [c['text'] for c in chunks]
    embeddings = embedder.encode(
        texts,
        batch_size=32,
        show_progress_bar=False
    )
    
    # Upload original file to S3
    try:
        s3.upload_file(
            filepath,
            S3_BUCKET,
            f'documents/{service}/{filename}'
        )
        REPORT['s3_uploaded'].append(
            f'documents/{service}/{filename}'
        )
    except Exception as e:
        print(f'S3 original upload failed: {e}')
    
    # Upload chunks JSON to S3
    chunks_data = {
        'filename': filename,
        'service': service,
        'total_chunks': len(chunks),
        'created': datetime.now(
            timezone.utc
        ).isoformat(),
        'chunks': chunks
    }
    try:
        s3.put_object(
            Bucket=S3_BUCKET,
            Key=(f'chunks/{service}/'
                 f'{filename}_chunks.json'),
            Body=json.dumps(
                chunks_data, indent=2
            ).encode('utf-8'),
            ContentType='application/json'
        )
        REPORT['s3_uploaded'].append(
            f'chunks/{service}/{filename}_chunks.json'
        )
    except Exception as e:
        print(f'S3 chunks upload failed: {e}')
    
    # Upload embeddings to S3
    try:
        emb_array = np.array(embeddings)
        with tempfile.NamedTemporaryFile(suffix='.npy', delete=False) as tmp:
            np.save(tmp.name, emb_array)
            s3.upload_file(
                tmp.name,
                S3_BUCKET,
                (f'embeddings/{service}/'
                 f'{filename}_embeddings.npy')
            )
            os.unlink(tmp.name)
        REPORT['s3_uploaded'].append(
            f'embeddings/{service}/'
            f'{filename}_embeddings.npy'
        )
    except Exception as e:
        print(f'S3 embeddings upload failed: {e}')
    
    # Add to FAISS index via bot API
    try:
        import requests
        txt_content = '\n\n'.join(texts)
        r = requests.post(
            'http://localhost:8000/'
            'api/v1/documents/upload',
            files={
                'file': (
                    filename + '.txt',
                    txt_content.encode('utf-8'),
                    'text/plain'
                )
            },
            timeout=120
        )
        if r.status_code == 200:
            print(f'Bot ingested: {filename}')
        else:
            print(
                f'Bot failed {filename}: '
                f'{r.status_code}'
            )
    except Exception as e:
        print(f'Bot ingest error: {e}')
    
    return len(chunks)


def process_single_file(filepath: str) -> dict:
    filename = Path(filepath).name
    
    # Detect AWS service from filename/path
    path_lower = filepath.lower()
    service = 'general'
    services = [
        'lambda','s3','ec2','iam','vpc',
        'rds','dynamodb','sagemaker','bedrock',
        'cloudwatch','ecs','eks','sns','sqs',
        'kinesis','glue','athena','redshift',
        'cloudformation','stepfunctions'
    ]
    for svc in services:
        if svc in path_lower:
            service = svc
            break
    
    # print(f'Processing: {filename} ({service})')
    
    # Extract text
    text = extract_text(filepath)
    
    if not text or len(text.strip()) < 5: # Lowered guard slightly for icons
        # print(f'No content: {filename}')
        return {'file': filename, 'status': 'skipped', 'reason': 'insufficient content'}
    
    # Chunk
    chunks = chunk_text(text, filename)
    
    # Vectorize and upload to S3
    n_chunks = vectorize_and_upload(
        filename, filepath, service, chunks
    )
    
    REPORT['files_processed'].append(filename)
    REPORT['total_chunks'] += n_chunks
    
    print(
        f'Done: {filename} '
        f'({n_chunks} chunks → S3)'
    )
    return {
        'file': filename,
        'status': 'success',
        'chunks': n_chunks,
        'service': service
    }


# ── WALK ALL FILES IN EXTRACTED DIR ────────────

all_files = []
for root, dirs, files in os.walk(EXTRACT_DIR):
    dirs[:] = [
        d for d in dirs 
        if not d.startswith('.')
        and d != '__MACOSX'
    ]
    for f in files:
        if f.startswith('.'):
            continue
        if f in ['Thumbs.db','desktop.ini']:
            continue
        # Check if extension is supported before adding to all_files
        ext = Path(f).suffix.lower()
        if ext in ('.pdf', '.png', '.jpg', '.jpeg', '.docx', '.pptx', '.xlsx', '.csv', '.txt', '.md', '.json'):
            all_files.append(os.path.join(root, f))

print(f'Total files found: {len(all_files)}')

# Limited processing
files_to_process = all_files[:MAX_FILES_TO_PROCESS]
print(f'Processing first {len(files_to_process)} files to demonstrate.')

# ── PARALLEL PROCESSING ─────────────────────────

results = []
with ThreadPoolExecutor(max_workers=4) as executor:
    futures = {
        executor.submit(
            process_single_file, fp
        ): fp for fp in files_to_process
    }
    for future in as_completed(futures):
        try:
            result = future.result()
            results.append(result)
        except Exception as e:
            fp = futures[future]
            print(f'Failed {fp}: {e}')
            REPORT['files_failed'].append(
                Path(fp).name
            )

# ── SAVE REPORT TO S3 ───────────────────────────

REPORT['finished'] = datetime.now().isoformat()
REPORT['summary'] = {
    'total_files_found': len(all_files),
    'limited_to': len(files_to_process),
    'processed': len(REPORT['files_processed']),
    'failed': len(REPORT['files_failed']),
    's3_objects': len(REPORT['s3_uploaded']),
    'total_chunks': REPORT['total_chunks']
}

report_json = json.dumps(REPORT, indent=2)
s3.put_object(
    Bucket=S3_BUCKET,
    Key='registry/training_report.json',
    Body=report_json.encode('utf-8'),
    ContentType='application/json'
)

print('\n' + '='*50)
print('PROCESSING COMPLETE')
print('='*50)
s = REPORT['summary']
print(f'Total Files Found: {s["total_files_found"]}')
print(f'Limited To:       {s["limited_to"]}')
print(f'Processed:        {s["processed"]}')
print(f'Failed:           {s["failed"]}')
print(f'S3 Objects:       {s["s3_objects"]}')
print(f'Total Chunks:     {s["total_chunks"]}')
print('Report saved to S3: registry/training_report.json')
